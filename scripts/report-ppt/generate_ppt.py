#!/usr/bin/env python3
"""
노출량 보고서 PPT 생성 (python-pptx).
stdin: ReportPptParams JSON
argv[1]: 출력 .pptx 경로
템플릿 이미지: TEMPLATES_DIR 또는 스크립트 기준 public/report-ppt-templates
"""
import json
import os
import sys
import base64
import tempfile
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 16:9
SLIDE_W = 10
SLIDE_H = 5.625

TEMPLATE_NAMES = {
    "cover": "slide-01-cover.png.jpg",
    "summary": "slide-02-summary.png.jpg",
    "banner": "slide-03-banner.png.jpg",
    "end": "slide-04-End.png.jpg",
}

# Windows에서 한글이 ??? 로 나오지 않도록 기본 한글 폰트 사용 (NanoSans KR 미설치 시)
FONT_NAME = "Malgun Gothic"  # Windows 맑은 고딕. Mac: "Apple SD Gothic Neo"


def _candidate_templates_dirs():
    """템플릿 폴더 후보: 환경변수 > 스크립트 기준 프로젝트 루트 > 현재 작업 디렉터리."""
    env = os.environ.get("REPORT_PPT_TEMPLATES_DIR")
    if env:
        yield Path(env)
    script_dir = Path(__file__).resolve().parent
    yield script_dir.parent / "public" / "report-ppt-templates"
    yield Path.cwd() / "public" / "report-ppt-templates"


def load_template_path(name: str) -> Optional[Path]:
    for templates_dir in _candidate_templates_dirs():
        if not templates_dir.is_dir():
            continue
        for fn in (TEMPLATE_NAMES[name], TEMPLATE_NAMES[name].replace(".png.jpg", ".png")):
            p = templates_dir / fn
            if p.exists():
                return p
    return None


def _find_pptx_template() -> Optional[Path]:
    """report_template.pptx 가 있으면 그 경로 반환 (덮어씌우기 대신 placeholder만 치환)."""
    for templates_dir in _candidate_templates_dirs():
        if not templates_dir.is_dir():
            continue
        p = templates_dir / "report_template.pptx"
        if p.exists():
            return p
    return None


def _replace_in_text(s: str, replacements: dict) -> str:
    for key, value in replacements.items():
        s = s.replace(key, str(value or ""))
    return s


def _replace_placeholders_in_shape(shape, replacements: dict) -> None:
    """placeholder 치환. 문단 전체 텍스트로 치환해 run이 쪼개져 있어도 동작."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            full = "".join(run.text for run in para.runs)
            if not full:
                continue
            new_text = _replace_in_text(full, replacements)
            if new_text == full:
                continue
            # 문단 내용을 치환된 텍스트 하나의 run으로 교체 (첫 run만 남기고 나머지 제거)
            for i, run in enumerate(para.runs):
                if i == 0:
                    run.text = new_text
                else:
                    run.text = ""
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text = _replace_in_text(cell.text, replacements)


def _replace_placeholders_in_presentation(prs: Presentation, replacements: dict) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_placeholders_in_shape(shape, replacements)


def format_man(n: float) -> str:
    if n >= 10_000:
        return f"{n / 10_000:.1f}만"
    return f"{n:,.0f}"


def format_number(n: float) -> str:
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}백만"
    if n >= 10_000:
        return f"{n / 10_000:.1f}만"
    return f"{n:,.0f}"


def add_background_image(slide, image_path: Path, prs_width, prs_height):
    """슬라이드 전체를 덮는 배경 이미지 추가. (먼저 추가해 맨 뒤 레이어에 두고, 이후 텍스트가 위에 그려짐.)"""
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0),
        Inches(0),
        width=Inches(SLIDE_W),
        height=Inches(SLIDE_H),
    )
    return pic


def _sanitize_text(s: str) -> str:
    """서로게이트 등 잘못된 유니코드 제거."""
    if not s:
        return s
    return "".join(c for c in s if not (0xD800 <= ord(c) <= 0xDFFF))


def _hex_to_rgb(hex_str: str):
    h = (hex_str or "000000").lstrip("#")[:6]
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color_hex=None, align=None):
    """color_hex: 템플릿 배경 위에는 "FFFFFF", 템플릿 없을 때는 None(진한색 333333)으로 가독성 확보."""
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _sanitize_text(text or "")
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = FONT_NAME
    hex_val = color_hex if color_hex is not None else "333333"
    p.font.color.rgb = _hex_to_rgb(hex_val)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    return tx


def main():
    if len(sys.argv) < 2:
        print("Usage: generate_ppt.py <output.pptx>", file=sys.stderr)
        sys.exit(1)
    out_path = Path(sys.argv[1])

    try:
        params = json.load(sys.stdin)
    except Exception as e:
        print(f"JSON parse error: {e}", file=sys.stderr)
        sys.exit(2)

    advertiser_name = _sanitize_text(params.get("advertiserName", "") or "")
    station = _sanitize_text(params.get("station", "") or "")
    line = _sanitize_text(params.get("line", "") or "")
    display_days = int(params.get("displayDays") or 0)
    exposure = params.get("exposure") or {}
    total_exposure = float(exposure.get("totalExposure") or 0)
    daily_flow = float(exposure.get("dailyFlow") or 0)
    by_time_band = exposure.get("byTimeBand") or []
    image_base64s = params.get("imageBase64s") or []
    subtitle = _sanitize_text(params.get("subtitle") or "")
    date_str = (params.get("dateStr") or "")[:8]
    campaign_manager_name = _sanitize_text(params.get("campaignManagerName") or "")
    campaign_manager_email = _sanitize_text(params.get("campaignManagerEmail") or "")

    date_fmt = f"{date_str[:4]}.{date_str[4:6]}.{date_str[6:8]}" if date_str and len(date_str) >= 8 else ""
    format_man_val = format_man(daily_flow) if daily_flow else "0"
    format_total = format_man(total_exposure) if total_exposure else "0"

    # .pptx 템플릿이 있으면 "덮어씌우기" 없이 placeholder만 치환 후 저장
    pptx_tpl = _find_pptx_template()
    if pptx_tpl:
        prs = Presentation(str(pptx_tpl))
        replacements = {
            "{광고주명}": advertiser_name,
            "{년.월.일}": date_fmt,
            "{캠페인 담당자 이름}": campaign_manager_name,
            "{캠페인 담당자 이메일}": campaign_manager_email,
            "{역/호선}": f"{line} {station}",
            "{역사명}": station,
            "{역명}": station,
            "{n호선}": line,
            "{5호선}": line,
            "{여의도역}": station,
            "{일평균 유동인구}": format_man_val,
            "{게재 기간}": f"{display_days}일",
            "{예상 총 노출량}": format_total,
            "{총 노출량}": format_total,
        }
        _replace_placeholders_in_presentation(prs, replacements)
        prs.save(str(out_path))
        return

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # ——— 슬라이드 1: 표지 ———
    s1 = prs.slides.add_slide(blank)
    t_cover = load_template_path("cover")
    if t_cover:
        add_background_image(s1, t_cover, SLIDE_W, SLIDE_H)
    add_textbox(s1, 0.8, 2.0, 8, 0.7, f"{advertiser_name} N.square 광고배너 게재보고서", font_size=24, bold=True, color_hex="FFFFFF")
    if date_str and len(date_str) >= 8:
        date_fmt = f"{date_str[:4]}.{date_str[4:6]}.{date_str[6:8]}"
        add_textbox(s1, 4, 5.0, 2, 0.35, date_fmt, font_size=14, color_hex="FFFFFF", align="center")
    if not t_cover:
        add_textbox(s1, 0.5, 1.2, 9, 0.8, f"{advertiser_name} 게재 현황 보고서", font_size=28, bold=True, align="center")
        add_textbox(s1, 0.5, 2.0, 9, 0.4, f"{line} {station}{' / ' + subtitle if subtitle else ''}", font_size=18, align="center")
        if date_str:
            add_textbox(s1, 0.5, 2.6, 9, 0.35, f"보고 일자: {date_str[:4]}-{date_str[4:6]}-{date_str[6:8]}", font_size=14, align="center")

    # ——— 슬라이드 2: 노출량 요약 ———
    s2 = prs.slides.add_slide(blank)
    t_summary = load_template_path("summary")
    if t_summary:
        add_background_image(s2, t_summary, SLIDE_W, SLIDE_H)
    by_tb = sorted(by_time_band, key=lambda x: -x.get("exposure", 0))
    peak = by_tb[0] if by_tb else None
    second = by_tb[1] if len(by_tb) > 1 else None
    add_textbox(s2, 0.5, 0.25, 6, 0.45, f"{advertiser_name} 예상 노출량 리포트", font_size=20, bold=True, color_hex="FFFFFF")
    add_textbox(s2, 7, 0.3, 2.2, 0.3, f"피크 시간대: {peak['band'] if peak else '-'} ({format_man(peak['exposure']) if peak else '0'}명)", font_size=11, color_hex="FFFFFF")
    if second:
        add_textbox(s2, 7, 0.55, 2.2, 0.3, f"2순위: {second['band']} ({format_man(second['exposure'])}명)", font_size=11, color_hex="FFFFFF")
    add_textbox(s2, 0.5, 0.95, 2, 0.35, f"{line} {station}", font_size=12, color_hex="FFFFFF")
    add_textbox(s2, 2.6, 0.95, 2, 0.35, f"{format_man(daily_flow)}명", font_size=12, color_hex="FFFFFF")
    add_textbox(s2, 4.7, 0.95, 1.2, 0.35, f"{display_days}일", font_size=12, color_hex="FFFFFF")
    add_textbox(s2, 6, 0.95, 2, 0.35, f"{format_man(total_exposure)}명", font_size=12, color_hex="FFFFFF")
    # 시간대별 표
    if by_tb:
        total_man = total_exposure / 10_000
        rows = 1 + len(by_tb)
        cols = 4
        table = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(4.2), Inches(0.35 * rows)).table
        table.cell(0, 0).text = "시간대"
        table.cell(0, 1).text = "예상 노출량(만명)"
        table.cell(0, 2).text = "비중"
        table.cell(0, 3).text = "순위"
        for i, b in enumerate(by_tb):
            exp = b.get("exposure", 0)
            pct = (exp / total_exposure * 100) if total_exposure else 0
            table.cell(i + 1, 0).text = f"{b.get('band', '')} 피크" if i == 0 else b.get("band", "")
            table.cell(i + 1, 1).text = f"{exp / 10_000:.1f}"
            table.cell(i + 1, 2).text = f"{pct:.0f}%"
            table.cell(i + 1, 3).text = str(i + 1)
    add_textbox(s2, 0.5, 5.15, 9, 0.4,
        "※ 본 수치는 공공데이터(역사별 승·하차/유동인구) 기반 추정치이며, 실제 광고 노출 수는 편차가 있을 수 있습니다.",
        font_size=8, color_hex="AAAAAA")
    if not t_summary:
        add_textbox(s2, 0.5, 0.3, 9, 0.5, "예상 노출량 요약", font_size=22, bold=True)
        summary_rows = [
            ["항목", "내용"],
            ["역/호선", f"{line} {station}"],
            ["일평균 유동인구", f"{format_number(daily_flow)}명"],
            ["게재 기간", f"{display_days}일"],
            ["예상 총 노출량", f"{format_number(total_exposure)}명"],
        ]
        t2 = s2.shapes.add_table(len(summary_rows), 2, Inches(0.5), Inches(1.0), Inches(9), Inches(0.4 * len(summary_rows))).table
        for r, row in enumerate(summary_rows):
            t2.cell(r, 0).text = row[0]
            t2.cell(r, 1).text = row[1]

    # ——— 슬라이드 3: 시간대별 (있을 때만) ———
    if by_time_band:
        s3 = prs.slides.add_slide(blank)
        add_textbox(s3, 0.5, 0.3, 9, 0.5, "시간대별 예상 노출량", font_size=22, bold=True)
        tr = [["시간대", "예상 노출량"]] + [[b.get("band", ""), f"{format_number(b.get('exposure', 0))}명"] for b in by_time_band]
        t3 = s3.shapes.add_table(len(tr), 2, Inches(0.5), Inches(1.0), Inches(9), Inches(0.4 * len(tr))).table
        for r, row in enumerate(tr):
            t3.cell(r, 0).text = row[0]
            t3.cell(r, 1).text = row[1]

    # ——— 슬라이드 4~: 촬영 광고 사진 (슬라이드당 2장) ———
    photo_w, photo_h = 4.5, 3.375
    for i in range(0, len(image_base64s), 2):
        pair = image_base64s[i : i + 2]
        slide = prs.slides.add_slide(blank)
        t_banner = load_template_path("banner")
        if t_banner:
            add_background_image(slide, t_banner, SLIDE_W, SLIDE_H)
        for idx, b64 in enumerate(pair):
            raw = b64.split(",", 1)[-1] if "," in b64 else b64
            try:
                img_data = base64.b64decode(raw)
            except Exception:
                continue
            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as f:
                f.write(img_data)
                tmp_path = f.name
            try:
                x = 0.4 if idx == 0 else 5.1
                slide.shapes.add_picture(tmp_path, Inches(x), Inches(1.2), width=Inches(photo_w), height=Inches(photo_h))
            finally:
                os.unlink(tmp_path)

    # ——— 마지막: 문서 끝 ———
    s_end = prs.slides.add_slide(blank)
    t_end = load_template_path("end")
    if t_end:
        add_background_image(s_end, t_end, SLIDE_W, SLIDE_H)
    add_textbox(s_end, 3, 2.4, 4, 0.5, f"케이티 나스미디어 {campaign_manager_name}", font_size=18, color_hex="FFFFFF", align="center")
    add_textbox(s_end, 3, 2.95, 4, 0.4, campaign_manager_email, font_size=14, color_hex="FFFFFF", align="center")
    if not t_end:
        add_textbox(s_end, 0.5, 2, 9, 0.5, "문서 끝", font_size=22, bold=True, align="center")
        if campaign_manager_name or campaign_manager_email:
            add_textbox(s_end, 0.5, 2.6, 9, 0.4, " · ".join(filter(None, [campaign_manager_name, campaign_manager_email])), font_size=14, align="center")

    prs.save(str(out_path))


if __name__ == "__main__":
    main()
